import sys
import os
sys.path.append(os.path.dirname(os.path.dirname(__file__)))
from pptx import Presentation
from pptx.util import Inches, Pt, Centipoints
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
import pdfTopptx.config as config
import pptx
# ContentObject := Text | Image | ...
# Assume that objectList is sorted by page & coordinates

from assembler.header import addHeader
from assembler.title import addFirstPage
# Version 1
def assemble(contentObjectList, outputFileName): # [] ContentObject
    prs = Presentation()
    addFirstPage(prs)
    titleSlideLayout = prs.slide_layouts[6] # blank
    currentPage = -1
    width = prs.slide_width
    height = prs.slide_height
    slideShape = None
    tf = None
    # flags
    isAddText = False
    isAddImage = False
    for contentObject in contentObjectList:
        if (currentPage != contentObject["pageNumber"]): #update page number
            slide = prs.slides.add_slide(titleSlideLayout)
            slideShape = slide.shapes 
            addHeader(slideShape, width, height / 10)
            currentPage += 1
            txBox = slideShape.add_textbox(width / 10, height / 6, 4 * width / 5, 0)
            tf = txBox.text_frame
            tf.word_wrap  = True
            isAddText = False
            isAddImage = False
            
            
        
        if contentObject["contentType"] == "text":
            isAddText = True
            for line in contentObject["content"].strip().replace('\n','@(br)').split('@(br)'):
                if (line != ''):
                    p = tf.add_paragraph()
                    p.text = ' '.join(line.split()).replace('.', '. ').strip()
                    p.level = 0
                    p.font.size = Pt(12)
                    p.font.name = config.FONT_FAMILY
                    p.font.size = Pt(config.PARAGRAPH_FONT_SIZE)
                    p.space_after = Pt(config.PARAGRAPH_FONT_SIZE)


        if contentObject["contentType"] == "image":
            isAddImage = True
            if (isAddText):
                slideShape.add_picture(contentObject["imagePath"], width / 4, 3 * height / 5, None, height / 5)
            else:
                slideShape.add_picture(contentObject["imagePath"], width / 4, height / 4, None, height / 2)
                
                
        if contentObject["contentType"] == "table":
            rows, columns = len(contentObject["dataFrame"]), len(contentObject["dataFrame"][0])
            x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
            # Add table to the slide
            graphic_frame = slide.shapes.add_table(rows, columns, x, y, cx, cy)
            table = graphic_frame.table

            # Populate the table with DataFrame values
            for i in range(rows):
                for j in range(columns):
                    cell = table.cell(i, j)
                    cell.text = str(contentObject["dataFrame"][i][j])

    prs.save(outputFileName)