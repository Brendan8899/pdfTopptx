from pptx import Presentation
import pdfTopptx.config  as config
prs = Presentation()


prs.save('test.pptx')

def addFirstPage(prs):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text_frame.paragraphs[0].font.name = config.FONT_FAMILY
    subtitle.text_frame.paragraphs[0].font.name = config.FONT_FAMILY # fix: not working
    
    title.text_frame.paragraphs[0].text = config.TITLE
    subtitle.text_frame.paragraphs[0].text = config.SUBTITLE