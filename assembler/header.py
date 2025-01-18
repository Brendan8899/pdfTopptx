import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
import config

def addHeader(slideShape, width, height):
    txBox = slideShape.add_textbox(0, 0, width, height)
    fill = txBox.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    text_frame = txBox.text_frame
    text_frame.margin_left = Pt(10)
    text_frame.margin_right = Pt(10)
    text_frame.margin_top = Pt(3)
    text_frame.margin_bottom = Pt(3)
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE 
    
    p = text_frame.add_paragraph()
    p.text = config.TITLE
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = config.FONT_FAMILY
    p.font.size = pptx.util.Pt(20)
    